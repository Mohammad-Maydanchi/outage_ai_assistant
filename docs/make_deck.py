"""Generate the demo deck: Vivant-Downtime-Check-Demo.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG     = RGBColor(0x0C, 0x10, 0x18)
CARD   = RGBColor(0x16, 0x1E, 0x2C)
TEXT   = RGBColor(0xE7, 0xEC, 0xF5)
MUTED  = RGBColor(0x9A, 0xA6, 0xBC)
ACCENT = RGBColor(0xF0, 0xA8, 0x3C)
GREEN  = RGBColor(0x6F, 0xD4, 0x9B)
MONO   = "Menlo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=18, color=TEXT, bold=False, first=False, align=PP_ALIGN.LEFT,
         font=None, space_after=8, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    runs = text if isinstance(text, list) else [(text, color, bold)]
    if bullet:
        r = p.add_run(); r.text = "•  "; r.font.color.rgb = ACCENT
        r.font.size = Pt(size); r.font.bold = True
        if font: r.font.name = font
    for item in runs:
        t, c, b = item if isinstance(item, tuple) else (item, color, bold)
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = c
        if font: r.font.name = font
    return p


def heading(s, title, kicker=None):
    tf = box(s, 0.7, 0.5, 12, 1.2)
    if kicker:
        para(tf, kicker.upper(), size=13, color=ACCENT, bold=True, first=True, space_after=2)
        para(tf, title, size=32, color=TEXT, bold=True)
    else:
        para(tf, title, size=32, color=TEXT, bold=True, first=True)
    # accent rule
    ln = s.shapes.add_shape(1, Inches(0.72), Inches(1.65), Inches(1.1), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def card(s, left, top, width, height, fill=CARD):
    sh = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = RGBColor(0x2A, 0x36, 0x4C); sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


# ---------- 1. Title ----------
s = slide()
tf = box(s, 0.9, 2.5, 11.5, 2.6)
para(tf, [("Vivant ", TEXT, True), ("Downtime Check", ACCENT, True)], size=54, first=True, space_after=6)
para(tf, "Outbound AI Voice Agent for ISP Outages", size=26, color=TEXT, space_after=18)
para(tf, "An AI agent that calls the ISP, navigates the phone system, and returns a clean, structured outage report — unattended.", size=18, color=MUTED)
tf2 = box(s, 0.9, 6.4, 11.5, 0.6)
para(tf2, "Proof-of-concept demo  ·  50 automated tests  ·  proven on live calls", size=14, color=MUTED, first=True)

# ---------- 2. The problem ----------
s = slide(); heading(s, "The problem", "Why this exists")
tf = box(s, 0.8, 2.1, 11.6, 4.5)
for t in [
    "MSPs and IT teams must call ISPs about outages across many client sites.",
    "They sit on hold, navigate phone menus, talk to reps or bots, then write up notes by hand.",
    "Slow, costly in human time, and inconsistent.",
]:
    para(tf, t, size=20, bullet=True, space_after=14)
para(tf, [("The differentiator is NOT the AI call (commoditized) — it's the ", MUTED, False),
          ("structured, ticketable report", ACCENT, True),
          (", produced unattended across many sites.", MUTED, False)], size=20, space_after=10)

# ---------- 3. What it does ----------
s = slide(); heading(s, "What it does", "The loop")
steps = [
    ("1.  Enter the outage", "A simple web form (stands in for a Utiliko / PSA ticket)."),
    ("2.  Dispatch the agent", "It calls the ISP and navigates whatever it reaches."),
    ("3.  Navigate & talk", "Human rep, IVR keypad menu, or a conversational AI bot."),
    ("4.  Write the report", "Claude reads the transcript — never fabricates."),
    ("5.  Read it on screen", "Outcome, summary, call recording, full transcript."),
]
top = 2.0
for title, sub in steps:
    card(s, 0.8, top, 11.7, 0.82)
    tf = box(s, 1.05, top + 0.08, 11.2, 0.7)
    para(tf, [(title + "    ", ACCENT, True), (sub, TEXT, False)], size=17, first=True)
    top += 0.95

# ---------- 4. Architecture ----------
s = slide(); heading(s, "How it's built", "Architecture")
tf = box(s, 0.8, 2.0, 11.6, 4.8)
rows = [
    ("Frontend", "React (Vite) — form, request list, report card."),
    ("Backend", "FastAPI + SQLite (thin data layer)."),
    ("Voice", "Vapi = telephony + speech-to-text + LLM + text-to-speech + turn-taking."),
    ("In-call brain", "GPT-4o-mini (fast)."),
    ("Report extraction", "Claude Sonnet (accurate, never-fabricate)."),
    ("Swappable", "VoiceProvider interface — Stub for tests, Vapi live, Bland could drop in."),
    ("Quality", "50 automated tests (TDD)."),
]
for k, v in rows:
    para(tf, [(f"{k:>18}   ", ACCENT, True), (v, TEXT, False)], size=17, bullet=False, space_after=11, font=None)

# ---------- 5. What's proven ----------
s = slide(); heading(s, "Proven end-to-end on live calls", "Test matrix")
data = [
    ("Live human rep", "PASS"),
    ("Simple keypad IVR menu", "PASS"),
    ("Hard branching IVR  (decoys + keypad data entry)", "PASS"),
    ("Conversational ISP bot  (self-service loop → rep)", "PASS"),
    ("Hard conversational bot  (verify gauntlet + double deflection + hold)", "PASS"),
]
top = 2.05
for label, status in data:
    card(s, 0.8, top, 11.7, 0.86)
    tf = box(s, 1.05, top + 0.16, 9.6, 0.6)
    para(tf, label, size=18, color=TEXT, first=True)
    tf2 = box(s, 10.7, top + 0.16, 1.6, 0.6)
    para(tf2, "✓ " + status, size=18, color=GREEN, bold=True, first=True, align=PP_ALIGN.CENTER)
    top += 0.99


def example_slide(kicker, title, transcript_lines, result):
    s = slide(); heading(s, title, kicker)
    card(s, 0.8, 2.0, 7.7, 4.9)
    tf = box(s, 1.0, 2.15, 7.3, 4.6)
    para(tf, "TRANSCRIPT (excerpt)", size=12, color=MUTED, bold=True, first=True, space_after=8)
    for who, line in transcript_lines:
        c = ACCENT if who == "Agent" else MUTED
        para(tf, [(f"{who}: ", c, True), (line, TEXT, False)], size=12.5, font=MONO, space_after=7)
    card(s, 8.75, 2.0, 3.75, 4.9, fill=RGBColor(0x12, 0x1A, 0x28))
    tf = box(s, 8.95, 2.2, 3.4, 4.5)
    para(tf, "STRUCTURED REPORT", size=12, color=MUTED, bold=True, first=True, space_after=10)
    for k, v in result:
        para(tf, k.upper(), size=11, color=MUTED, space_after=1)
        para(tf, v, size=15, color=TEXT, bold=True, space_after=10)


# ---------- 6. Example: Hard IVR ----------
example_slide(
    "Example 1", "Hard IVR — keypad navigation",
    [("Agent", "I'm calling about a business internet outage."),
     ("ISP", "For technical support, press 3."),
     ("Agent", "[presses 3]"),
     ("ISP", "For internet, press 2."),
     ("Agent", "[presses 2]"),
     ("ISP", "To report an outage, press 2."),
     ("Agent", "[presses 2]"),
     ("ISP", "Enter your 10-digit number followed by #."),
     ("Agent", "[enters number + #]"),
     ("ISP", "Outage due to severe weather. ETA 9 PM. Ref 67890.")],
    [("Outcome", "Outage confirmed — ETA"), ("Reason", "Severe weather"),
     ("ETA", "9 PM today"), ("Reference #", "67890"), ("Responder", "Recorded message")],
)

# ---------- 7. Example: Conversational bot ----------
example_slide(
    "Example 2", "Conversational ISP bot",
    [("ISP", "How can I help you today?"),
     ("Agent", "I'm calling about a business internet outage at..."),
     ("ISP", "Can I send you a text with more info?"),
     ("Agent", "No thank you — I'd prefer to get it over the phone."),
     ("ISP", "(repeats the offer)"),
     ("Agent", "I'd like to speak with a representative, please."),
     ("ISP", "One moment, connecting you to tech support."),
     ("ISP", "Hi, this is technical support. How can I help?"),
     ("Agent", "Is there a known outage? Reason, ETA, ticket #?"),
     ("ISP", "Fiber cut. ETA 7 PM. Ticket 99887.")],
    [("Outcome", "Outage confirmed — ETA"), ("Reason", "Fiber cut"),
     ("ETA", "7 PM today"), ("Reference #", "99887"), ("Responder", "Human (rep)")],
)

# ---------- 8. Example: Hard conversational gauntlet ----------
example_slide(
    "Example 3  ·  the showcase", "Hard conversational — the gauntlet",
    [("ISP", "Which service — internet, TV, or phone?"),
     ("Agent", "Internet service."),
     ("ISP", "Can I text you the link?  (x2)"),
     ("Agent", "No thank you — connect me to a rep, please."),
     ("ISP", "First, verify: account number?"),
     ("Agent", "1234567890."),
     ("ISP", "PIN?"), ("Agent", "4321."),
     ("ISP", "Confirm the address?"),
     ("Agent", "2727 Linden, Dallas 75234."),
     ("ISP", "Hold or callback?"), ("Agent", "I'll hold."),
     ("ISP", "Damaged fiber line. ETA 8 PM. Ticket 55432.")],
    [("Outcome", "Outage confirmed — ETA"), ("Reason", "Damaged fiber line"),
     ("ETA", "8 PM today"), ("Reference #", "55432"), ("Responder", "Human (rep)")],
)

# ---------- 8b. Listen — embedded recording ----------
DOCS = "/Users/mohammadmaydanchi/Documents/vivant_corp/docs"
# build a poster image for the audio
from PIL import Image, ImageDraw
poster = Image.new("RGB", (1280, 720), (18, 26, 40))
dr = ImageDraw.Draw(poster)
dr.ellipse([540, 250, 740, 450], outline=(240, 168, 60), width=8)
dr.polygon([(605, 305), (605, 395), (685, 350)], fill=(240, 168, 60))
poster.save(f"{DOCS}/audio-poster.png")

s = slide(); heading(s, "Listen — the hard gauntlet call", "Real recording")
s.shapes.add_movie(f"{DOCS}/gauntlet-recording.wav", Inches(0.9), Inches(2.2),
                   Inches(5.6), Inches(3.15), poster_frame_image=f"{DOCS}/audio-poster.png",
                   mime_type="audio/x-wav")
tf = box(s, 6.9, 2.3, 5.6, 4.4)
para(tf, "~3.5 minutes  ·  your agent vs. a hard AT&T-style bot.", size=20, color=TEXT, first=True, space_after=18)
for t in ["Pushes past a double \"can I text you?\" deflection.",
          "Passes the verification gauntlet (account #, PIN, address).",
          "Holds for a rep, then extracts the outage facts.",
          "Hangs up cleanly — no loop."]:
    para(tf, t, size=16, color=MUTED, bullet=True, space_after=11)
para(tf, "(Press play in PowerPoint / Keynote.)", size=13, color=ACCENT, space_after=4)

# ---------- 8c. Full transcript ----------
s = slide(); heading(s, "Full transcript — the gauntlet call", "Complete record")
lines = [l for l in open(f"{DOCS}/full_transcript.txt").read().splitlines() if l.strip()]
half = (len(lines) + 1) // 2
for ci, chunk in enumerate((lines[:half], lines[half:])):
    tf = box(s, 0.8 + ci * 6.15, 1.95, 6.0, 5.2)
    for i, line in enumerate(chunk):
        if line.startswith("AI:"):
            who, rest, c = "Agent:", line[3:], ACCENT
        elif line.startswith("User:"):
            who, rest, c = "ISP:", line[5:], MUTED
        else:
            who, rest, c = "", line, TEXT
        para(tf, [(who + " ", c, True), (rest.strip(), TEXT, False)],
             size=10.5, font=MONO, first=(i == 0), space_after=5)

# ---------- 9. Never fabricate ----------
s = slide(); heading(s, "The golden rule: never fabricate", "Trust")
tf = box(s, 0.8, 2.1, 11.6, 4.5)
for t in [
    "Claude extracts only what was actually said in the call.",
    "If the ETA, reason, or ticket number was not given → it reports \"not provided\".",
    "If the call is too unclear to judge → it flags \"needs review\" for a human.",
    "It never invents an ETA or a ticket number to look complete.",
]:
    para(tf, t, size=20, bullet=True, space_after=16)
para(tf, [("The honesty IS the product", ACCENT, True),
          (" — a faithful, ticketable report an MSP can trust.", MUTED, False)], size=20)

# ---------- 10. Engineering rigor ----------
s = slide(); heading(s, "Stress-tested — bugs found & fixed", "Engineering rigor")
tf = box(s, 0.8, 2.0, 11.6, 4.8)
for t in [
    "Start-collision: two AIs greeting at once → one now waits.",
    "Goodbye-loop: agent wouldn't hang up → fixed end-call phrases.",
    "AI-to-AI turn-taking: tuned with LiveKit smart endpointing.",
    "Graceful errors: clean message when a call can't start (not a crash).",
    "Daily call limit on Vapi numbers → imported a Twilio number for outbound (no cap).",
]:
    para(tf, t, size=18, bullet=True, space_after=14)

# ---------- 11. Key insight ----------
s = slide(); heading(s, "Key insight: ISPs are becoming AI", "Strategic")
tf = box(s, 0.8, 2.2, 11.6, 4.2)
para(tf, "AT&T's real phone system is already a conversational AI.", size=22, color=TEXT, first=True, space_after=16)
para(tf, [("So calling an ISP is increasingly ", MUTED, False),
          ("our AI talking to their AI", ACCENT, True), (".", MUTED, False)], size=22, space_after=16)
para(tf, "AI-to-AI turn-taking is a real, ongoing challenge — and the reason we test "
        "against realistic conversational bots, not just menus.", size=20, color=MUTED)

# ---------- 12. Future works ----------
s = slide(); heading(s, "What's next", "Roadmap")
tf = box(s, 0.8, 2.0, 11.6, 4.8)
for t in [
    "Compliance for real-ISP calls (AI disclosure, recording consent).",
    "Branded caller ID — remove the \"Spam Likely\" label.",
    "Per-ISP provider routing (use the best voice platform per ISP).",
    "LLM-as-a-judge — auto-score each call's quality.",
    "Bulk multi-site sweeps; report sharing / export.",
]:
    para(tf, t, size=18, bullet=True, space_after=13)

# ---------- 13. Status ----------
s = slide()
tf = box(s, 0.9, 2.6, 11.5, 2.6)
para(tf, "Status", size=16, color=ACCENT, bold=True, first=True, space_after=8)
para(tf, "Core POC complete, tested, and proven on every ISP call type.", size=30, color=TEXT, bold=True, space_after=18)
para(tf, "50 automated tests passing  ·  code on GitHub  ·  real live-call demos available on request.",
     size=18, color=MUTED)

prs.save("/Users/mohammadmaydanchi/Documents/vivant_corp/docs/Vivant-Downtime-Check-Demo.pptx")
print("Saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
